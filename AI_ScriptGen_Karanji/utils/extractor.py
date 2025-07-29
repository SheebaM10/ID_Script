import os
from pptx import Presentation
from docx import Document
from pdfminer.high_level import extract_text as extract_pdf_text

def extract_text_from_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.pptx':
            return extract_from_pptx(filepath)
        elif ext == '.docx':
            return extract_from_docx(filepath)
        elif ext == '.pdf':
            return extract_from_pdf(filepath)
        else:
            return f"[Extraction Error] Unsupported file type: {ext}"
    except Exception as e:
        return f"[Extraction Error] {str(e)}"

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image
from io import BytesIO

def run_ocr(image_bytes):
    try:
        img = Image.open(BytesIO(image_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        return f"[OCR Error] {str(e)}"

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from PIL import Image
from io import BytesIO

def run_ocr(image_bytes):
    try:
        img = Image.open(BytesIO(image_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        return f"[OCR Error] {str(e)}"

def extract_from_pptx(path):
    try:
        prs = Presentation(path)
        slides = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_lines = [f"--- Slide {i} ---"]

            # Slide content
            for shape in slide.shapes:
                # Heading + bullet points
                if hasattr(shape, "text") and shape.text.strip():
                    lines = [line.strip() for line in shape.text.strip().splitlines() if line.strip()]
                    if lines:
                        heading = lines[0]
                        points = lines[1:]
                        slide_lines.append(f"\n{heading}")
                        for point in points:
                            slide_lines.append(f"- {point}")

                # Table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_lines.append("- " + " | ".join(row_text))

                # Image OCR
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    ocr_text = run_ocr(image_bytes)
                    if ocr_text:
                        slide_lines.append(f"[IMAGE CONTENT]: {ocr_text}")

            # Notes content
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
                if note_text:
                    slide_lines.append(f"[NOTES]: {note_text}")

            slides.append("\n".join(slide_lines))

        return "\n\n".join(slides).strip()

    except Exception as e:
        return f"[PPTX Extraction Error] {str(e)}"




from docx import Document

def extract_from_docx(path):
    try:
        doc = Document(path)
        paragraphs = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())

        # Extract and OCR images
        image_dir = "temp_images"
        os.makedirs(image_dir, exist_ok=True)
        image_texts = []

        for rel in doc.part._rels:
            if "image" in doc.part._rels[rel].target_ref:
                image_data = doc.part._rels[rel]._target.stream.read()
                image_path = os.path.join(image_dir, f"img_{len(image_texts)}.png")
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_data)

                # Run OCR on image
                ocr_result = run_ocr(image_path)
                image_texts.append(f"[IMAGE CONTENT]: {ocr_result}")

        paragraphs.extend(image_texts)

        return "\n\n".join(paragraphs)
    except Exception as e:
        return f"[DOCX Extraction Error] {str(e)}"




def extract_from_pdf(path):
    try:
        text = extract_pdf_text(path)
        return text.strip()
    except Exception as e:
        return f"[PDF Extraction Error] {str(e)}"

import pytesseract
from PIL import Image

def run_ocr(image_path):
    try:
        img = Image.open(image_path)
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        return f"[OCR Error] {str(e)}"
